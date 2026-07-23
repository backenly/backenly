# -*- coding: utf-8 -*-
"""
Backenly investor pitch deck generator.
Dark + white theme, restrained violet-500 (#8B5CF6) accents.
Output: Backenly-Pitch-Deck.pptx (16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0E, 0x0E, 0x10)   # deep near-black
PANEL    = RGBColor(0x17, 0x17, 0x1B)   # card
PANEL2   = RGBColor(0x1B, 0x1B, 0x20)   # card alt
VPANEL   = RGBColor(0x1A, 0x15, 0x26)   # violet-tinted card
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ZINC100  = RGBColor(0xF4, 0xF4, 0xF5)
ZINC300  = RGBColor(0xD4, 0xD4, 0xD8)
ZINC400  = RGBColor(0xA1, 0xA1, 0xAA)
ZINC500  = RGBColor(0x80, 0x80, 0x8A)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)   # violet-500 (brand)
VIOLET_L = RGBColor(0xA7, 0x8B, 0xFA)   # violet-400
VIOLET_D = RGBColor(0x5B, 0x3D, 0xA6)   # darker for fills
BORDER   = RGBColor(0x2A, 0x2A, 0x30)
GREEN    = RGBColor(0x4A, 0xDE, 0x80)

FONT   = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_S = "Segoe UI Semibold"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def txt(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of dicts {runs:[(text,{opts})], align, space_after, line_spacing, bullet}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ): pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if "space_after" in p: para.space_after = Pt(p["space_after"])
        if "space_before" in p: para.space_before = Pt(p["space_before"])
        if "line_spacing" in p: para.line_spacing = p["line_spacing"]
        runs = p["runs"] if isinstance(p["runs"], list) else [(p["runs"], {})]
        for t, o in runs:
            r = para.add_run(); r.text = t
            f = r.font
            f.name = o.get("font", FONT)
            f.size = Pt(o.get("size", 16))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", ZINC300)
            if o.get("italic"): f.italic = True
            sp = o.get("spacing")
            if sp is not None:
                rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(sp))
    return tb


def kicker(s, label, x=0.75, y=0.62):
    rect(s, x, y + 0.02, 0.22, 0.22, fill=VIOLET, rounded=True, radius=0.45)  # dot
    txt(s, x + 0.34, y - 0.06, 8, 0.4,
        [{"runs": [(label.upper(), {"size": 12.5, "bold": True, "color": VIOLET_L,
                                    "font": FONT_S, "spacing": 220})]}])


def title(s, text, x=0.75, y=0.95, w=11.8, size=33):
    txt(s, x, y, w, 1.2,
        [{"runs": [(text, {"size": size, "bold": True, "color": WHITE, "font": FONT_S})],
          "line_spacing": 1.0}])


def footer(s, n):
    rect(s, 0, 7.18, 13.333, 0.012, fill=BORDER)
    txt(s, 0.75, 7.2, 4, 0.3,
        [{"runs": [("Backenly", {"size": 9.5, "bold": True, "color": ZINC400}),
                   (" •", {"size": 9.5, "bold": True, "color": VIOLET})]}])
    txt(s, 8.0, 7.2, 4.583, 0.3,
        [{"runs": [("Confidential — Pre-seed 2026", {"size": 9.5, "color": ZINC500})],
          "align": PP_ALIGN.RIGHT}])
    txt(s, 12.4, 7.2, 0.6, 0.3,
        [{"runs": [(f"{n:02d}", {"size": 9.5, "bold": True, "color": VIOLET_L})],
          "align": PP_ALIGN.RIGHT}])


def card(s, x, y, w, h, fill=PANEL, line=BORDER):
    return rect(s, x, y, w, h, fill=fill, line=line, line_w=1.0, rounded=True, radius=0.07)


def check(s, x, y, text, size=13, color=ZINC300):
    txt(s, x, y, 0.4, 0.35, [{"runs": [("✓", {"size": size, "bold": True, "color": GREEN})]}])
    txt(s, x + 0.34, y, 9, 0.5, [{"runs": [(text, {"size": size, "color": color})]}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
# faint violet glow panel top-right
rect(s, 9.3, -1.2, 5.5, 4.0, fill=VPANEL, rounded=True, radius=0.2)
rect(s, 9.55, -1.0, 5.0, 3.6, fill=BG, rounded=True, radius=0.2)
# wordmark
txt(s, 0.85, 1.5, 8, 0.7,
    [{"runs": [("Backenly", {"size": 30, "bold": True, "color": WHITE, "font": FONT_S}),
               ("•", {"size": 30, "bold": True, "color": VIOLET})]}])
# headline
txt(s, 0.82, 2.5, 11.2, 2.0,
    [{"runs": [("Describe what you want to build.\n", {"size": 46, "bold": True, "color": WHITE, "font": FONT_S})],
      "line_spacing": 1.02},
     {"runs": [("It becomes a real ", {"size": 46, "bold": True, "color": WHITE, "font": FONT_S}),
               ("backend.", {"size": 46, "bold": True, "color": VIOLET_L, "font": FONT_S})],
      "line_spacing": 1.02}])
# sub
txt(s, 0.85, 4.65, 10.6, 0.9,
    [{"runs": [("Paste your idea from ChatGPT → a complete, live backend → connect it to "
                "Cursor, Bolt or Lovable in one line. No code, no SQL, no servers.",
                {"size": 16.5, "color": ZINC400})], "line_spacing": 1.25}])
# raise chip
chip = rect(s, 0.85, 5.7, 3.05, 0.5, fill=VPANEL, line=VIOLET_D, line_w=1.0, rounded=True, radius=0.5)
txt(s, 0.85, 5.72, 3.05, 0.46,
    [{"runs": [("Raising $100K  ·  Pre-seed", {"size": 13, "bold": True, "color": VIOLET_L})],
      "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
# founders / contact footer
rect(s, 0.85, 6.62, 11.6, 0.012, fill=BORDER)
txt(s, 0.85, 6.78, 11.6, 0.5,
    [{"runs": [("Adarsh Chiriyamkandath Jose", {"size": 12.5, "bold": True, "color": ZINC100}),
               ("  &  ", {"size": 12.5, "color": ZINC500}),
               ("Lakshmi Koonath", {"size": 12.5, "bold": True, "color": ZINC100}),
               ("      backenly.com   ·   @Backenly   ·   hello@backenly.com",
                {"size": 12, "color": ZINC400})]}])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Problem")
title(s, "AI broke the frontend open. The backend is still a wall.")
probs = [
    ("The frontend is generated — the backend still needs an engineer.",
     "Tables, APIs, auth, RLS, storage, realtime, deploy, monitoring: each a specialist skill."),
    ("BaaS relocated the problem; it didn't solve it.",
     "Supabase & Firebase hand you tools. You still design schemas and operate it. Built for engineers."),
    ("AI code generators make it worse over time.",
     "Bolt / Base44 give you backend code you now own, debug, deploy and maintain."),
    ("Nobody runs it for you.",
     "It breaks at 2am — a permissive policy, an error spike, schema drift — and no one is watching."),
]
y = 1.95
for i, (h, b) in enumerate(probs):
    cy = y + i * 1.02
    card(s, 0.75, cy, 11.83, 0.9)
    rect(s, 0.75, cy, 0.07, 0.9, fill=VIOLET, rounded=False)
    txt(s, 1.05, cy + 0.13, 11.3, 0.35, [{"runs": [(h, {"size": 15.5, "bold": True, "color": WHITE})]}])
    txt(s, 1.05, cy + 0.5, 11.3, 0.35, [{"runs": [(b, {"size": 12.5, "color": ZINC400})]}])
# punchline
pl = rect(s, 0.75, 6.18, 11.83, 0.72, fill=VPANEL, line=VIOLET_D, rounded=True, radius=0.12)
txt(s, 1.05, 6.18, 11.3, 0.72,
    [{"runs": [("A population 10× larger than backend engineers can ship a UI — ",
                {"size": 15, "color": ZINC300}),
               ("but not a product.", {"size": 15, "bold": True, "color": VIOLET_L})]}],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Solution")
title(s, "One sentence → a live production backend that runs itself.")
txt(s, 0.75, 1.62, 11.8, 0.6,
    [{"runs": [("A single agentic AI brain grounds itself in your live backend, plans, and executes "
                "~80 governed tools to produce a running system — not code you maintain.",
                {"size": 14, "color": ZINC400})], "line_spacing": 1.2}])
pillars = [
    ("The AI is the infrastructure",
     "Not a code printer. No repo to maintain. Every project runs in its own isolated PostgreSQL schema — true database-level multi-tenancy."),
    ("Safe by construction",
     "Every change is dry-run, atomic, audited and confirm-gated. The AI asks before guessing and refuses raw SQL or manual schema edits by design."),
    ("It runs itself",
     "A live closed-loop autonomy runtime monitors health and self-heals drift — company-funded, so it never spends the user's credits."),
    ("Lives where builders already are",
     "A 1-line SDK for any frontend, plus a published MCP server so Cursor & Claude Code drive the backend in plain English."),
]
gx, gy, gw, gh = 0.75, 2.5, 5.83, 1.95
for i, (h, b) in enumerate(pillars):
    cx = gx + (i % 2) * (gw + 0.17)
    cy = gy + (i // 2) * (gh + 0.2)
    card(s, cx, cy, gw, gh, fill=PANEL)
    txt(s, cx + 0.35, cy + 0.28, 0.6, 0.5,
        [{"runs": [(f"0{i+1}", {"size": 20, "bold": True, "color": VIOLET})]}])
    txt(s, cx + 0.35, cy + 0.72, gw - 0.7, 0.4,
        [{"runs": [(h, {"size": 17, "bold": True, "color": WHITE})]}])
    txt(s, cx + 0.35, cy + 1.14, gw - 0.7, 0.7,
        [{"runs": [(b, {"size": 12.5, "color": ZINC400})], "line_spacing": 1.18}])
footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUCT MAP
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Product — fully built & live")
title(s, "A complete backend platform. Every section shipped.")
txt(s, 0.75, 1.6, 11.8, 0.5,
    [{"runs": [("Not a prototype — running in production at backenly.com today. "
                "~80 governed AI tools span every surface below.", {"size": 13.5, "color": ZINC400})]}])
groups = [
    ("DATA", ["Tables — NL → real Postgres", "APIs — auto REST + aggregates"]),
    ("ACCESS", ["Auth — JWT, OAuth, Teams/B2B", "Storage — S3 buckets, signed URLs"]),
    ("CONNECT", ["Frontend SDK — 1-line client", "MCP — Cursor / Claude Code"]),
    ("DEPLOY", ["Publish — live API URL", "Versioned deploys + rollback"]),
    ("SERVICES", ["Integrations — Stripe, AI, email", "Realtime + Functions + RAG"]),
    ("OPERATIONS", ["Monitoring — metrics + anomalies", "Autonomy — self-healing loop"]),
]
gx, gy, gw, gh = 0.75, 2.25, 3.83, 1.95
for i, (label, items) in enumerate(groups):
    cx = gx + (i % 3) * (gw + 0.17)
    cy = gy + (i // 3) * (gh + 0.2)
    card(s, cx, cy, gw, gh)
    rect(s, cx + 0.3, cy + 0.32, 0.16, 0.16, fill=VIOLET, rounded=True, radius=0.45)
    txt(s, cx + 0.56, cy + 0.24, gw - 0.8, 0.35,
        [{"runs": [(label, {"size": 12.5, "bold": True, "color": VIOLET_L, "spacing": 180, "font": FONT_S})]}])
    for j, it in enumerate(items):
        txt(s, cx + 0.32, cy + 0.78 + j * 0.5, gw - 0.6, 0.45,
            [{"runs": [(it, {"size": 13, "color": ZINC300})], "line_spacing": 1.05}])
footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEMO (screenshot placeholders)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Demo")
title(s, "Type in plain English. Watch your backend build itself.")
# big main screenshot slot
def shot(x, y, w, h, label, sub):
    r = rect(s, x, y, w, h, fill=PANEL2, line=VIOLET_D, line_w=1.25, rounded=True, radius=0.04)
    txt(s, x, y + h/2 - 0.5, w, 0.4,
        [{"runs": [("◳  SCREENSHOT", {"size": 12, "bold": True, "color": VIOLET_L, "spacing": 150})],
          "align": PP_ALIGN.CENTER}])
    txt(s, x + 0.2, y + h/2 - 0.08, w - 0.4, 0.5,
        [{"runs": [(label, {"size": 14.5, "bold": True, "color": ZINC100})], "align": PP_ALIGN.CENTER}])
    txt(s, x + 0.2, y + h/2 + 0.32, w - 0.4, 0.5,
        [{"runs": [(sub, {"size": 11, "color": ZINC500})], "align": PP_ALIGN.CENTER}])
shot(0.75, 1.95, 7.4, 4.9, "AI chat building a backend",
     "Inspector — chat + live build timeline")
shot(8.32, 1.95, 4.26, 2.35, "Database / Tables view",
     "Tables + governed schema")
shot(8.32, 4.5, 4.26, 2.35, "Autonomy / Trust panel",
     "Self-healing + approvals")
footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE MOAT (Autonomy)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "The moat")
title(s, "A backend that keeps itself alive — funded by us.")
txt(s, 0.75, 1.65, 7.3, 3.2,
    [{"runs": [("A live closed-loop runtime continuously compares your backend to a desired "
                "state and self-heals drift — missing keys, over-permissive policies, broken "
                "integrations — each root-caused, not patched.", {"size": 14.5, "color": ZINC300})],
      "line_spacing": 1.3, "space_after": 14},
     {"runs": [("Hard safety contract:", {"size": 13.5, "bold": True, "color": WHITE})], "space_after": 6}])
safety = [
    "Per-project circuit breaker that fails closed",
    "Dial: Off / Conservative / Balanced / Aggressive",
    "Incident change-freeze when health is anomalous",
    "Plain-English “because” for every action + audit",
    "Auth / destructive changes are NEVER auto-applied",
]
for i, it in enumerate(safety):
    check(s, 0.78, 3.55 + i * 0.46, it, size=12.5)
# right: funded card
card(s, 8.35, 1.95, 4.23, 4.6, fill=VPANEL, line=VIOLET_D)
txt(s, 8.7, 2.35, 3.55, 2.2,
    [{"runs": [("“Your credits are for ", {"size": 19, "color": ZINC300, "font": FONT_S}),
               ("building", {"size": 19, "bold": True, "color": WHITE, "font": FONT_S}),
               (". Keeping it alive is ", {"size": 19, "color": ZINC300, "font": FONT_S}),
               ("on us", {"size": 19, "bold": True, "color": VIOLET_L, "font": FONT_S}),
               (".”", {"size": 19, "color": ZINC300, "font": FONT_S})], "line_spacing": 1.15}])
rect(s, 8.7, 5.0, 3.5, 0.012, fill=VIOLET_D)
txt(s, 8.7, 5.2, 3.55, 1.2,
    [{"runs": [("The self-healing loop is company-funded on paid tiers — a moat tied to our "
                "cost structure, not a feature competitors can toggle on.",
                {"size": 12, "color": ZINC400})], "line_spacing": 1.25}])
footer(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MARKET
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Market")
title(s, "Taking BaaS into the 10×-larger AI-builder segment.")
mk = [
    ("TAM", "$95B+", "Global BaaS + backend / cloud-app-platform spend, expanding as software creation shifts to natural language.", PANEL),
    ("SAM", "$15–25B", "Indie devs, pre-PMF→growth startups, agencies, and the AI-coding ecosystem (Supabase / Firebase / Xano / Railway slice).", PANEL),
    ("SOM", "Low 8-figure ARR", "The AI-assisted builder segment that can't use existing BaaS without an engineer. $19–$99 ACV.", VPANEL),
]
gx, gw = 0.75, 3.83
for i, (lab, big, body, fill) in enumerate(mk):
    cx = gx + i * (gw + 0.17)
    card(s, cx, 2.2, gw, 3.6, fill=fill, line=VIOLET_D if fill == VPANEL else BORDER)
    txt(s, cx + 0.4, 2.5, gw - 0.8, 0.4,
        [{"runs": [(lab, {"size": 13, "bold": True, "color": VIOLET_L, "spacing": 220, "font": FONT_S})]}])
    txt(s, cx + 0.4, 3.0, gw - 0.8, 1.0,
        [{"runs": [(big, {"size": 30, "bold": True, "color": WHITE, "font": FONT_S})], "line_spacing": 0.95}])
    txt(s, cx + 0.4, 4.25, gw - 0.8, 1.4,
        [{"runs": [(body, {"size": 12.5, "color": ZINC400})], "line_spacing": 1.25}])
txt(s, 0.75, 6.1, 11.8, 0.6,
    [{"runs": [("Primary buyers:  ", {"size": 13.5, "bold": True, "color": WHITE}),
               ("AI-assisted developers · agent-native builders · frontend-only developers · "
                "indie hackers · founders shipping with AI tools.", {"size": 13.5, "color": ZINC400})]}])
footer(s, 7)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BUSINESS MODEL
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Business model")
title(s, "Flat-rate SaaS, metered by active users — not nickel-and-dimed.")
# pricing table
rows = [
    ["Plan", "Price / mo", "Active users", "AI credits", "Autonomy"],
    ["Free — Sandbox", "$0", "500 MAU", "200 / mo", "24h · bounded"],
    ["Starter — Builder", "$19", "50,000 MAU", "2,000 / mo", "2h · funded"],
    ["Pro — Scale", "$99", "250,000 MAU", "6,000 / mo", "30m · funded"],
    ["Enterprise", "Custom", "Custom", "Custom", "+ RBAC · SSO · SLA"],
]
tx, ty, tw, th = 0.75, 2.05, 11.83, 3.0
tbl = s.shapes.add_table(len(rows), len(rows[0]), Inches(tx), Inches(ty), Inches(tw), Inches(th)).table
tbl.first_row = False; tbl.horz_banding = False
widths = [3.2, 1.7, 2.5, 2.0, 2.43]
for c, wv in enumerate(widths):
    tbl.columns[c].width = Inches(wv)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.6)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.18); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = VIOLET_D
        elif rows[ri][0].startswith("Pro"):
            cell.fill.solid(); cell.fill.fore_color.rgb = VPANEL
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if ri % 2 else PANEL2
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.name = FONT_S if (ri == 0 or ci == 0) else FONT
        r.font.size = Pt(13 if ri == 0 else 13.5)
        r.font.bold = (ri == 0 or ci == 0)
        r.font.color.rgb = WHITE if ri == 0 else (ZINC100 if ci == 0 else ZINC300)
# credits note
card(s, 0.75, 5.35, 11.83, 1.15, fill=PANEL)
txt(s, 1.05, 5.5, 11.3, 0.9,
    [{"runs": [("Token-backed AI credits  ", {"size": 14, "bold": True, "color": VIOLET_L}),
               ("(1 credit = 1,000 tokens — published & stable). ", {"size": 13.5, "color": ZINC300}),
               ("Cost scales with real work, so the “paste a whole spec for 1 message” exploit is "
                "structurally gone. We never meter tables, endpoints, rows or deploys — only AI, "
                "MAU, storage & realtime.", {"size": 13.5, "color": ZINC400})], "line_spacing": 1.2}])
footer(s, 8)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GO TO MARKET
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Go-to-market")
title(s, "Be the default backend for AI-built frontends.")
txt(s, 0.75, 1.6, 11.8, 0.5,
    [{"runs": [("The wedge: go where AI-assisted developers already are — and live "
                "inside the AI coding tools themselves.", {"size": 14, "color": ZINC400})]}])
gtm = [
    ("Launch channels", ["Product Hunt — live build-sim demo", "X / @Backenly — build-in-public", "LinkedIn — founder narrative"]),
    ("Distribution wedge", ["Default backend for v0 / Bolt / Lovable", "Native MCP citizen in Cursor & Claude Code", "“Connect in one line” inside the tool"]),
    ("Content & community", ["“Shipped a SaaS, no backend engineer”", "Indie Hackers · Framer · Webflow", "Supabase-setup vs. one-sentence demos"]),
]
gx, gw = 0.75, 3.83
for i, (h, items) in enumerate(gtm):
    cx = gx + i * (gw + 0.17)
    card(s, cx, 2.4, gw, 3.5)
    txt(s, cx + 0.32, 2.7, gw - 0.6, 0.4,
        [{"runs": [(h, {"size": 15.5, "bold": True, "color": WHITE})]}])
    rect(s, cx + 0.34, 3.2, 0.5, 0.022, fill=VIOLET)
    for j, it in enumerate(items):
        rect(s, cx + 0.34, 3.5 + j * 0.7 + 0.07, 0.13, 0.13, fill=VIOLET, rounded=True, radius=0.45)
        txt(s, cx + 0.6, 3.5 + j * 0.7, gw - 0.85, 0.62,
            [{"runs": [(it, {"size": 12.5, "color": ZINC300})], "line_spacing": 1.1}])
txt(s, 0.75, 6.25, 11.8, 0.5,
    [{"runs": [("Mistake to avoid:  ", {"size": 13, "bold": True, "color": VIOLET_L}),
               ("don’t chase engineers who already build backends — target builders with no alternative.",
                {"size": 13, "color": ZINC400})]}])
footer(s, 9)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITION
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Competition")
title(s, "Others help you build. We build it — and run it.")
rows = [
    ["", "Backenly", "Supabase", "Firebase", "Xano", "Base44 / Bolt"],
    ["NL → DB + REST API", "✓", "—", "—", "—", "—"],
    ["Governed atomic exec + rollback", "✓", "—", "—", "—", "—"],
    ["NL RLS, triggers, functions, RAG", "✓", "Manual", "Cloud fn", "—", "—"],
    ["AI operates a running system", "✓", "—", "—", "—", "—"],
    ["Self-healing autonomy (funded)", "✓", "—", "—", "—", "—"],
    ["Native MCP for Cursor / Claude Code", "✓", "—", "—", "—", "—"],
    ["Real relational Postgres you own", "✓", "✓", "—", "✓", "n/a"],
]
tx, ty, tw, th = 0.75, 2.0, 11.83, 4.0
tbl = s.shapes.add_table(len(rows), len(rows[0]), Inches(tx), Inches(ty), Inches(tw), Inches(th)).table
tbl.first_row = False; tbl.horz_banding = False
cw = [4.0, 1.75, 1.52, 1.52, 1.2, 1.84]
for c, wv in enumerate(cw):
    tbl.columns[c].width = Inches(wv)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.5)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        is_backenly = (ci == 1)
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = VIOLET_D if is_backenly else PANEL2
        elif is_backenly:
            cell.fill.solid(); cell.fill.fore_color.rgb = VPANEL
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL if ri % 2 else PANEL2
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.name = FONT
        r.font.size = Pt(12.5 if ci == 0 else 14)
        r.font.bold = (ri == 0 or ci == 0 or (val == "✓" and is_backenly))
        if val == "✓":
            r.font.color.rgb = WHITE if is_backenly else GREEN
        elif val == "—":
            r.font.color.rgb = ZINC500
        else:
            r.font.color.rgb = WHITE if (ri == 0 or ci == 0) else ZINC300
footer(s, 10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TRACTION
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Traction / milestones")
title(s, "Users = 0. Product = done. We’re launching today.")
txt(s, 0.75, 1.62, 11.8, 0.5,
    [{"runs": [("The 18 months & $1M+ a typical seed startup spends building is already live — "
                "and running itself.", {"size": 14, "color": ZINC400})]}])
# left: shipped checklist
card(s, 0.75, 2.35, 6.7, 4.2)
txt(s, 1.05, 2.6, 6, 0.4, [{"runs": [("Shipped to production", {"size": 15, "bold": True, "color": WHITE})]}])
ships = [
    "Agentic brain — ~80 governed tools + vision",
    "Every section live (Data → Autonomy)",
    "Funded self-healing autonomy loop",
    "RAG, Teams/B2B, push, webhooks, cron",
    "Published @backenly/mcp-server",
    "Per-project Postgres isolation + audit",
    "Token-backed credit billing (Paddle)",
]
for i, it in enumerate(ships):
    check(s, 1.05, 3.15 + i * 0.46, it, size=12.5)
# right: milestones timeline
card(s, 7.62, 2.35, 4.96, 4.2, fill=VPANEL, line=VIOLET_D)
txt(s, 7.92, 2.6, 4.4, 0.4, [{"runs": [("Roadmap", {"size": 15, "bold": True, "color": WHITE})]}])
miles = [
    ("Now", "Public launch — PH / X / LinkedIn. First paid cohort via founder’s pricing."),
    ("0–3 mo", "Eval harness + hardening as real users surface edges. First case studies."),
    ("3–9 mo", "MCP / AI-tool partnerships, agent-ecosystem word-of-mouth, first meaningful MRR."),
]
for i, (t, b) in enumerate(miles):
    yy = 3.25 + i * 1.05
    rect(s, 7.92, yy + 0.05, 0.14, 0.14, fill=VIOLET, rounded=True, radius=0.45)
    txt(s, 8.16, yy - 0.04, 4.2, 0.35, [{"runs": [(t, {"size": 13.5, "bold": True, "color": VIOLET_L})]}])
    txt(s, 8.16, yy + 0.3, 4.2, 0.7, [{"runs": [(b, {"size": 11.5, "color": ZINC400})], "line_spacing": 1.15}])
footer(s, 11)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TEAM
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Team")
title(s, "A rare builder-operator pair.")
team = [
    ("Adarsh Chiriyamkandath Jose", "Founder · Engineering",
     "Architected and built the entire platform: the agentic brain and its ~80-tool governance "
     "kernel, multi-tenant Postgres isolation, the funded self-healing runtime, the SDK, the MCP "
     "server and the production infra. Shipped the scope of a full team — and operates it live."),
    ("Lakshmi Koonath", "Co-founder · Product & GTM",
     "Owns the category narrative (“the AI is the backend”), pricing strategy (token-backed credits "
     "+ funded autonomy) and the launch motion targeting frontend-first and non-technical builders."),
]
for i, (name, role, bio) in enumerate(team):
    cx = 0.75 + i * 6.0
    card(s, cx, 2.3, 5.83, 4.1)
    rect(s, cx + 0.4, 2.7, 0.9, 0.9, fill=VPANEL, line=VIOLET_D, rounded=True, radius=0.25)
    initials = "".join([w[0] for w in name.split()[:2]]).upper()
    txt(s, cx + 0.4, 2.7, 0.9, 0.9,
        [{"runs": [(initials, {"size": 22, "bold": True, "color": VIOLET_L, "font": FONT_S})],
          "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx + 0.4, 3.78, 5.0, 0.4, [{"runs": [(name, {"size": 18, "bold": True, "color": WHITE, "font": FONT_S})]}])
    txt(s, cx + 0.4, 4.2, 5.0, 0.35, [{"runs": [(role, {"size": 12.5, "bold": True, "color": VIOLET_L})]}])
    txt(s, cx + 0.4, 4.65, 5.05, 1.6, [{"runs": [(bio, {"size": 12.5, "color": ZINC400})], "line_spacing": 1.25}])
footer(s, 12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THE ASK
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "The ask")
title(s, "Raising $100K to turn a finished product into users.")
txt(s, 0.75, 1.65, 7.2, 1.0,
    [{"runs": [("The product is built and live — we’re not raising to build it. "
                "We’re raising to ", {"size": 15, "color": ZINC300}),
               ("distribute it and operate it at scale.", {"size": 15, "bold": True, "color": WHITE})],
      "line_spacing": 1.3}])
# big number
card(s, 8.35, 1.9, 4.23, 1.85, fill=VPANEL, line=VIOLET_D)
txt(s, 8.35, 2.05, 4.23, 0.9,
    [{"runs": [("$100K", {"size": 46, "bold": True, "color": WHITE, "font": FONT_S})], "align": PP_ALIGN.CENTER}])
txt(s, 8.35, 3.05, 4.23, 0.5,
    [{"runs": [("Pre-seed  ·  ~12-month runway", {"size": 13, "color": ZINC400})], "align": PP_ALIGN.CENTER}])
# use of funds bars
funds = [
    ("Growth & GTM", 40, "Launch, content, MCP/AI-tool partnerships, agent-ecosystem channels"),
    ("Engineering", 30, "1 hire; eval harness, abuse limits, deeper autonomy"),
    ("Infra & AI cost", 20, "Postgres/compute + company-funded autonomy spend"),
    ("Ops & buffer", 10, "Runway buffer & operations"),
]
y = 3.0
for i, (lab, pct, desc) in enumerate(funds):
    yy = y + i * 0.92
    txt(s, 0.78, yy, 3.0, 0.35, [{"runs": [(lab, {"size": 14, "bold": True, "color": WHITE})]}])
    txt(s, 0.78, yy + 0.36, 6.5, 0.35, [{"runs": [(desc, {"size": 11.5, "color": ZINC500})]}])
    # bar track
    rect(s, 3.6, yy + 0.02, 3.6, 0.26, fill=PANEL2, rounded=True, radius=0.5)
    rect(s, 3.6, yy + 0.02, 3.6 * pct / 100.0, 0.26, fill=VIOLET, rounded=True, radius=0.5)
    txt(s, 7.3, yy - 0.02, 0.8, 0.35, [{"runs": [(f"{pct}%", {"size": 14, "bold": True, "color": VIOLET_L})]}])
footer(s, 13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — APPENDIX: TARGET CUSTOMERS
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Appendix — who we sell to")
title(s, "Target customers, by fit.")
tiers = [
    ("TIER 1 — BEST FIT", VIOLET_L, [
        ("Frontend-only developers", "★★★★★"),
        ("AI-assisted developers", "★★★★★"),
        ("AI application builders", "★★★★★"),
        ("Indie hackers", "★★★★☆"),
        ("Hackathon teams", "★★★★☆"),
        ("Non-technical founders", "★★★★☆"),
    ]),
    ("TIER 2 — GOOD WITH FRICTION", ZINC300, [
        ("SaaS builders (Teams/RBAC live)", "★★★★☆"),
        ("Early-stage founders", "★★★★☆"),
        ("Agencies (MVP delivery)", "★★★☆☆"),
    ]),
    ("TIER 3 — POOR FIT", ZINC500, [
        ("Full-stack devs (want control)", "—"),
        ("Enterprise needing self-hosting", "—"),
        ("Pure no-code (SDK blocker)", "—"),
        ("Analytics-heavy (raw SQL)", "—"),
    ]),
]
gx, gw = 0.75, 3.83
for i, (head, col, items) in enumerate(tiers):
    cx = gx + i * (gw + 0.17)
    card(s, cx, 2.2, gw, 4.4)
    txt(s, cx + 0.32, 2.45, gw - 0.6, 0.4,
        [{"runs": [(head, {"size": 12, "bold": True, "color": col, "spacing": 120, "font": FONT_S})]}])
    rect(s, cx + 0.34, 2.95, gw - 0.68, 0.018, fill=BORDER)
    for j, (nm, st) in enumerate(items):
        yy = 3.15 + j * 0.55
        txt(s, cx + 0.32, yy, gw - 1.2, 0.4, [{"runs": [(nm, {"size": 12.5, "color": ZINC300})]}])
        txt(s, cx + gw - 1.05, yy, 0.85, 0.4,
            [{"runs": [(st, {"size": 11.5, "color": VIOLET if st.startswith("★") else ZINC500})],
              "align": PP_ALIGN.RIGHT}])
footer(s, 14)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 9.3, 4.6, 5.5, 4.0, fill=VPANEL, rounded=True, radius=0.2)
rect(s, 9.05, 4.85, 5.0, 3.6, fill=BG, rounded=True, radius=0.2)
txt(s, 0.85, 2.15, 8, 0.6,
    [{"runs": [("Backenly", {"size": 22, "bold": True, "color": WHITE, "font": FONT_S}),
               ("•", {"size": 22, "bold": True, "color": VIOLET})]}])
txt(s, 0.82, 2.9, 11.4, 1.8,
    [{"runs": [("Give a real backend to everyone who can already\nbuild the ",
                {"size": 36, "bold": True, "color": WHITE, "font": FONT_S}),
               ("front", {"size": 36, "bold": True, "color": VIOLET_L, "font": FONT_S}),
               (" of their dream.", {"size": 36, "bold": True, "color": WHITE, "font": FONT_S})],
      "line_spacing": 1.05}])
rect(s, 0.85, 4.95, 2.2, 0.03, fill=VIOLET)
txt(s, 0.85, 5.3, 11, 0.5,
    [{"runs": [("Adarsh Chiriyamkandath Jose  ·  Lakshmi Koonath",
                {"size": 14, "bold": True, "color": ZINC100})]}])
txt(s, 0.85, 5.75, 11, 0.5,
    [{"runs": [("backenly.com    ·    @Backenly    ·    hello@backenly.com",
                {"size": 13, "color": ZINC400})]}])

prs.save("Backenly-Pitch-Deck.pptx")
print("Saved Backenly-Pitch-Deck.pptx with", len(prs.slides._sldIdLst), "slides")
